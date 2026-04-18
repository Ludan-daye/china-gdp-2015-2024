"""生成项目展示 PPT —— FT 编辑部风，16:9，与 HTML/PNG 视觉一致。

共 22 张幻灯片：
  1  封面
  2  关键数据摘要
  3  目录
  4  Ch.1 章节页
  5-9  任务 1 图 1-4 + H1 英雄图
  10  Ch.2 章节页
  11-16  任务 2 图 1-5 + H2 英雄图
  17  Ch.3 章节页
  18-20  任务 3 图 1-3
  21  结论四条
  22  方法论与鸣谢
"""
from __future__ import annotations
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PROJECT_ROOT = Path(__file__).resolve().parent.parent
OUTPUTS_DIR = PROJECT_ROOT / "outputs"
SLIDES_DIR = PROJECT_ROOT / "slides"
SLIDES_DIR.mkdir(exist_ok=True)

# 色板（同 _common.py）
C_BG = RGBColor(0xF3, 0xEE, 0xE0)
C_BG_DEEP = RGBColor(0xEA, 0xE3, 0xD0)
C_INK = RGBColor(0x1B, 0x1B, 0x1B)
C_INK_SOFT = RGBColor(0x3A, 0x3A, 0x3A)
C_MUTED = RGBColor(0x76, 0x69, 0x5A)
C_RULE = RGBColor(0xC8, 0xBE, 0xA8)
C_TERTIARY = RGBColor(0xE0, 0x7B, 0x39)
C_SECONDARY = RGBColor(0x3A, 0x6F, 0xA0)
C_HIGHLIGHT = RGBColor(0xB2, 0x32, 0x2A)
C_PRIMARY = RGBColor(0x7A, 0x8B, 0x5E)

FONT_SERIF = "Songti SC"
FONT_SANS = "PingFang SC"
FONT_MONO = "Menlo"


def fill_bg(slide, color=C_BG):
    """为整张幻灯片填充背景色。"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                slide.part.package.presentation_part.presentation.slide_width,
                                slide.part.package.presentation_part.presentation.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    # 置底
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def txt(slide, left, top, width, height, text, *,
        font=FONT_SANS, size=14, bold=False, italic=False,
        color=C_INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        line_spacing=1.2):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def hline(slide, left, top, width, color=C_INK, weight=1.4):
    ln = slide.shapes.add_connector(1, left, top, left + width, top)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def add_topbar(slide, current=None):
    """所有内容页顶部小导航条。"""
    slide_w = Inches(13.333)
    txt(slide, Inches(0.45), Inches(0.22), Inches(6), Inches(0.3),
        "■  工业 → 服务 · 中国经济的十年",
        font=FONT_SERIF, size=11, bold=True, color=C_INK)
    # 右侧导航
    nav_items = [
        ("Ch.1 趋势", "ch1"),
        ("Ch.2 区域", "ch2"),
        ("Ch.3 相关", "ch3"),
        ("结语", "coda"),
    ]
    x = Inches(8.2)
    for label, tag in nav_items:
        color = C_TERTIARY if tag == current else C_MUTED
        bold = (tag == current)
        t = txt(slide, x, Inches(0.22), Inches(1.2), Inches(0.3),
                label, font=FONT_MONO, size=9,
                bold=bold, color=color, align=PP_ALIGN.LEFT)
        x += Inches(1.15)
    # 水平细分隔线
    hline(slide, Inches(0.45), Inches(0.55), Inches(12.4),
          color=C_RULE, weight=0.5)


def add_source_footer(slide):
    txt(slide, Inches(0.45), Inches(7.12), Inches(12), Inches(0.25),
        "数据来源：国家统计局《2024 年国民经济和社会发展统计公报》及各年度分省数据",
        font=FONT_MONO, size=8, color=C_MUTED)


# ----------------------------------------------------------------------
# SLIDE BUILDERS
# ----------------------------------------------------------------------

def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill_bg(s)
    # 左上 kicker
    txt(s, Inches(0.7), Inches(0.7), Inches(6), Inches(0.4),
        "—— 期末项目 · 数据可视化 · 2026",
        font=FONT_MONO, size=10, color=C_HIGHLIGHT)
    # 主标题（两行）
    txt(s, Inches(0.7), Inches(1.6), Inches(12), Inches(1.6),
        "从工业驱动",
        font=FONT_SERIF, size=72, bold=True, color=C_INK)
    txt(s, Inches(0.7), Inches(2.9), Inches(12), Inches(1.6),
        "走向  服务驱动",
        font=FONT_SERIF, size=72, bold=True, color=C_INK)
    # "服务驱动" 暖橙高亮块（橙色矩形）
    hl = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            Inches(3.35), Inches(3.85), Inches(3.4), Inches(0.28))
    hl.fill.solid(); hl.fill.fore_color.rgb = C_TERTIARY
    hl.line.fill.background()
    # 覆盖上一行"服务驱动"文字（改为橙色描边）——再写一层橙色 "服务驱动"
    txt(s, Inches(3.35), Inches(2.9), Inches(5), Inches(1.6),
        "服务驱动",
        font=FONT_SERIF, size=72, bold=True, color=C_TERTIARY)
    # 需要手动重写左侧 "走向  "（因为上一个覆盖了）—— 不重复，因为"走向"没被遮挡；
    # 实际上橙色文字要位于"走向  "右侧。左侧"走向"已经在上面那行，不受影响。

    # 副标题
    txt(s, Inches(0.7), Inches(4.6), Inches(11), Inches(1.4),
        "以 2015–2024 十年数据为镜，讲述中国经济规模翻倍、\n产业结构换骨、区域格局重塑的故事。",
        font=FONT_SERIF, size=22, color=C_INK_SOFT, line_spacing=1.35)
    # 底部 meta 栏
    hline(s, Inches(0.7), Inches(6.5), Inches(12), color=C_INK, weight=1.5)
    meta = [
        ("时段", "2015 — 2024"),
        ("样本", "31 省 · 全国"),
        ("叙事", "结构转型 A 线"),
        ("交付", "14 图 · 报告 · PPT"),
    ]
    for i, (k, v) in enumerate(meta):
        x = Inches(0.7 + i * 3.0)
        txt(s, x, Inches(6.65), Inches(2.8), Inches(0.25), k,
            font=FONT_MONO, size=8.5, color=C_MUTED)
        txt(s, x, Inches(6.90), Inches(2.8), Inches(0.3), v,
            font=FONT_SERIF, size=13, bold=True, color=C_INK)


def slide_summary(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_topbar(s)
    txt(s, Inches(0.7), Inches(0.95), Inches(11), Inches(0.6),
        "三个关键数字，十年中国经济浓缩",
        font=FONT_SERIF, size=32, bold=True, color=C_INK)
    hline(s, Inches(0.7), Inches(1.65), Inches(12), color=C_INK, weight=2)

    # 三栏大数字
    cards = [
        ("70.3  →  134.8", "万亿元", "全国 GDP 总量", "十年翻倍"),
        ("51.7  →  56.8", "% · 占 GDP 比重", "第三产业占比", "服务化加速"),
        ("+0.79", "Pearson r · 省际增速", "31 省增速同步性", "经济一盘棋"),
    ]
    col_w = Inches(4.0)
    for i, (num, unit, label, tag) in enumerate(cards):
        x = Inches(0.7 + i * 4.15)
        # 大数字
        txt(s, x, Inches(2.1), col_w, Inches(1.6), num,
            font=FONT_SERIF, size=54, bold=True, color=C_INK)
        # 单位
        txt(s, x, Inches(3.6), col_w, Inches(0.4), unit,
            font=FONT_SERIF, size=13, italic=True, color=C_MUTED)
        # 横线
        hline(s, x, Inches(4.35), Inches(3.8), color=C_RULE, weight=0.8)
        # 标签
        txt(s, x, Inches(4.5), col_w, Inches(0.5), label,
            font=FONT_SANS, size=16, bold=True, color=C_INK_SOFT)
        # tag
        txt(s, x, Inches(5.0), col_w, Inches(0.5), "· " + tag,
            font=FONT_SERIF, size=15, italic=True, color=C_TERTIARY)

    add_source_footer(s)


def slide_toc(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_topbar(s)
    txt(s, Inches(0.7), Inches(0.95), Inches(11), Inches(0.6),
        "全篇由三章组成",
        font=FONT_SERIF, size=32, bold=True, color=C_INK)
    hline(s, Inches(0.7), Inches(1.65), Inches(12), color=C_INK, weight=2)

    chapters = [
        ("01", "总量在膨胀，结构在换骨",
         "5 张图 · 全国 GDP 走势 / 三产堆叠 / 占比演化 / 增速节奏 / H1 三元图"),
        ("02", "四极稳定，梯队重排",
         "6 张图 · 排名柱 / 三产占比 / 哑铃图 / 分层地图 / 小多图 / H2 结构时钟"),
        ("03", "经济是联动的，也是异质的",
         "3 张图 · Treemap / 结构-增速散点 / 省际增速相关热图"),
    ]
    for i, (no, title, sub) in enumerate(chapters):
        top = Inches(2.2 + i * 1.45)
        # 大号 roman
        txt(s, Inches(0.7), top, Inches(1.6), Inches(1.3), no,
            font=FONT_SERIF, size=54, italic=True, color=C_TERTIARY)
        # 章节名
        txt(s, Inches(2.4), top + Inches(0.1), Inches(10), Inches(0.7), title,
            font=FONT_SERIF, size=26, bold=True, color=C_INK)
        # sub
        txt(s, Inches(2.4), top + Inches(0.85), Inches(10), Inches(0.5), sub,
            font=FONT_SANS, size=13, color=C_INK_SOFT)
        # 细分割
        if i < 2:
            hline(s, Inches(0.7), top + Inches(1.32), Inches(12),
                  color=C_RULE, weight=0.5)


def slide_chapter_divider(prs, no: str, roman: str, title: str, sub: str, current: str):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C_BG_DEEP)
    # 不画顶栏（章节页需要留白）
    # 大字 roman 居左
    txt(s, Inches(0.8), Inches(1.6), Inches(3), Inches(0.5),
        f"CHAPTER · {no}", font=FONT_MONO, size=13,
        color=C_MUTED, bold=True)
    txt(s, Inches(0.8), Inches(2.2), Inches(3), Inches(4), roman,
        font=FONT_SERIF, size=180, italic=True, color=C_TERTIARY,
        line_spacing=0.9)
    # 右侧大字标题
    txt(s, Inches(4.5), Inches(2.8), Inches(8.2), Inches(2),
        title, font=FONT_SERIF, size=50, bold=True, color=C_INK,
        line_spacing=1.1)
    # 副标
    txt(s, Inches(4.5), Inches(4.6), Inches(8.2), Inches(1.6),
        sub, font=FONT_SERIF, size=18, italic=True, color=C_INK_SOFT,
        line_spacing=1.5)
    # 底部横线 + chapter count
    hline(s, Inches(0.8), Inches(6.75), Inches(11.7), color=C_INK, weight=1.5)


def slide_chart(prs, png: str, kicker: str, title: str, sub: str,
                source_hint: str = "", hero: bool = False, current: str = "ch1"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_topbar(s, current=current)

    # Kicker（任务编号）
    kcolor = C_TERTIARY
    tb = txt(s, Inches(0.7), Inches(0.78), Inches(6), Inches(0.3),
             kicker, font=FONT_MONO, size=9, bold=True, color=kcolor)
    # HERO 标签
    if hero:
        hero_tag = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(11.4), Inches(0.75), Inches(1.4), Inches(0.3))
        hero_tag.fill.solid(); hero_tag.fill.fore_color.rgb = C_TERTIARY
        hero_tag.line.fill.background()
        tf = hero_tag.text_frame; tf.margin_left = Emu(0); tf.margin_right = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = "HERO IMAGE"
        run.font.name = FONT_MONO; run.font.size = Pt(9); run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    # 标题
    txt(s, Inches(0.7), Inches(1.08), Inches(12), Inches(0.55),
        title, font=FONT_SERIF, size=22, bold=True, color=C_INK)
    # 副标
    txt(s, Inches(0.7), Inches(1.62), Inches(12), Inches(0.5),
        sub, font=FONT_SERIF, size=12.5, italic=True, color=C_INK_SOFT)

    # 细分隔线
    hline(s, Inches(0.7), Inches(2.15), Inches(12), color=C_RULE, weight=0.5)

    # 图片 —— 按短边限高插入，保持宽高比
    img_path = OUTPUTS_DIR / png
    pic = s.shapes.add_picture(str(img_path),
                               Inches(0.7), Inches(2.3),
                               height=Inches(4.6))
    # 居中
    slide_w = Inches(13.333)
    pic.left = int((slide_w - pic.width) / 2)

    # 来源+图表类型脚注（两端对齐）
    txt(s, Inches(0.7), Inches(7.05), Inches(6), Inches(0.3),
        "数据来源：国家统计局", font=FONT_MONO, size=8, color=C_MUTED)
    if source_hint:
        txt(s, Inches(7), Inches(7.05), Inches(6), Inches(0.3),
            source_hint, font=FONT_MONO, size=8, color=C_MUTED,
            align=PP_ALIGN.RIGHT)


def slide_conclusion(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s, C_BG_DEEP)
    add_topbar(s, current="coda")
    # 标题
    txt(s, Inches(0.7), Inches(0.95), Inches(12), Inches(0.6),
        "结语 · Coda", font=FONT_MONO, size=11,
        bold=True, color=C_HIGHLIGHT)
    txt(s, Inches(0.7), Inches(1.35), Inches(12), Inches(1.2),
        "量变的表皮，质变的里子。",
        font=FONT_SERIF, size=40, bold=True, color=C_INK)
    hline(s, Inches(0.7), Inches(2.55), Inches(12), color=C_INK, weight=2)

    # 4 条 takeaway
    items = [
        ("01", "总量翻倍，但增速在换挡",
         "十年名义 GDP 从 70 万亿到 134.8 万亿，接近翻一番。\n2020 年疫情短暂拐点后，增速从高速切换至中速区间。"),
        ("02", "服务化是时代主旋律",
         "31 省全部向第三产业顶点漂移。\n北京三产占比 85%、上海 79%，已进入典型服务型经济区间。"),
        ("03", "四极稳固，但梯队在重排",
         "粤苏鲁浙仍是四极，合计 34.9%。\n西部省份累计增速更快；山西为唯一 2024 较 2023 小幅收缩的省份。"),
        ("04", "增速同步，路径却分岔",
         "省际增速 r = +0.79 高度联动；\n但结构升级 ≠ 快速增长 (r = −0.21)，追赶与升级不等价。"),
    ]
    for i, (no, head, body) in enumerate(items):
        col = i % 2; row = i // 2
        x = Inches(0.7 + col * 6.15)
        y = Inches(2.95 + row * 2.05)
        # 编号
        txt(s, x, y, Inches(5.9), Inches(0.5), no,
            font=FONT_SERIF, size=24, italic=True, color=C_TERTIARY)
        # 标题
        txt(s, x, y + Inches(0.5), Inches(5.9), Inches(0.5), head,
            font=FONT_SERIF, size=18, bold=True, color=C_INK)
        # 正文
        txt(s, x, y + Inches(1.0), Inches(5.9), Inches(1), body,
            font=FONT_SERIF, size=11.5, color=C_INK_SOFT, line_spacing=1.45)
        # 顶部横线（与编号平齐上方）
        hline(s, x, y - Inches(0.06), Inches(5.9), color=C_INK, weight=1.2)


def slide_colophon(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fill_bg(s)
    add_topbar(s)
    txt(s, Inches(0.7), Inches(0.95), Inches(4), Inches(0.5),
        "方法论  ·  METHODOLOGY", font=FONT_MONO, size=12,
        bold=True, color=C_HIGHLIGHT)
    txt(s, Inches(0.7), Inches(1.45), Inches(12), Inches(0.8),
        "数据、工具与复现方式",
        font=FONT_SERIF, size=30, bold=True, color=C_INK)
    hline(s, Inches(0.7), Inches(2.25), Inches(12), color=C_INK, weight=2)

    rows = [
        ("数据时段", "2015 – 2024（年度）"),
        ("数据口径", "国家统计局第五次全国经济普查（2023）修订后数据"),
        ("覆盖范围", "中国大陆 31 省 / 自治区 / 直辖市（港澳台未纳入）"),
        ("主数据源", "NBS《2024 年国民经济和社会发展统计公报》及历年分省数据"),
        ("地图底图", "阿里云 DataV · 中国省级 GeoJSON"),
        ("绘图工具", "Python · pandas · numpy · matplotlib · squarify（无 geopandas）"),
        ("视觉体系", "FT 编辑部风 · 暖橙 = 第三产业 · 雾蓝 = 第二产业 · 低饱和绿 = 第一产业"),
        ("数据文件", "data/gdp_2015_2024.csv（长表 310 行）· gdp_wide_total.csv · gdp_national.csv"),
        ("交付件", "14 张 PNG · 项目报告 report.md · HTML 站点 · 本 PPT"),
        ("在线查看", "https://ludan-daye.github.io/china-gdp-2015-2024/"),
    ]
    y = Inches(2.55)
    for k, v in rows:
        txt(s, Inches(0.7), y, Inches(2.6), Inches(0.35), k,
            font=FONT_MONO, size=9.5, color=C_MUTED, bold=True)
        txt(s, Inches(3.4), y, Inches(9.5), Inches(0.35), v,
            font=FONT_SANS, size=12, color=C_INK)
        hline(s, Inches(0.7), y + Inches(0.38), Inches(12.1),
              color=C_RULE, weight=0.3)
        y += Inches(0.42)


# ----------------------------------------------------------------------
def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ===== 开篇 =====
    slide_cover(prs)
    slide_summary(prs)
    slide_toc(prs)

    # ===== Ch.1 总体趋势 =====
    slide_chapter_divider(prs, "01", "I",
                          "总量在膨胀，\n结构在换骨。",
                          "Chapter 01 · 总体经济发展趋势 · 5 张图（含 H1 三元图）",
                          "ch1")
    slide_chart(prs, "task1_national_gdp.png",
                "T1.1  ·  全国 GDP 走势",
                "十年翻倍：2015–2024 全国 GDP 走势",
                "名义 GDP 从 70.3 万亿增至 134.8 万亿，十年间规模接近翻一番。",
                "折线图 · Line Chart", current="ch1")
    slide_chart(prs, "task1_industry_stack.png",
                "T1.2  ·  三次产业堆叠面积",
                "结构换骨：三次产业增加值的十年积累",
                "暖橙色第三产业面积从 36.3 万亿长至 76.6 万亿，增量远超二产。",
                "Stacked Area · 堆叠面积图", current="ch1")
    slide_chart(prs, "task1_industry_share.png",
                "T1.3  ·  三次产业占比演化",
                "服务化加速：第三产业占比从 51.7% 升至 56.8%",
                "十年间三产占比提升 5 个百分点，二产退让 3.5 个百分点。",
                "Line Chart · 占比折线", current="ch1")
    slide_chart(prs, "task1_growth_rate.png",
                "T1.4  ·  年度同比增速",
                "V 型恢复：名义 GDP 增速在疫情后显著反弹",
                "2020 年跌至 2.9%，2021 年反弹至 13.4%，此后回归中速区间。",
                "Bar Chart · 柱状图", current="ch1")
    slide_chart(prs, "task1_h1_ternary.png",
                "H1  ·  封面级 HERO 图",
                "漂向第三产业顶点：31 省十年产业结构轨迹",
                "灰色 = 31 省轨迹，红色 = 4 个极端代表。所有箭头方向一致。",
                "Ternary Plot · 三元图", hero=True, current="ch1")

    # ===== Ch.2 区域对比 =====
    slide_chapter_divider(prs, "02", "II",
                          "四极稳定，\n梯队重排。",
                          "Chapter 02 · 地区经济规模对比 · 6 张图（含 H2 结构时钟）",
                          "ch2")
    slide_chart(prs, "task2_ranking.png",
                "T2.1  ·  2024 年 GDP 排名",
                "四极格局稳固：粤苏鲁浙贡献全国 1/3 GDP",
                "广东、江苏双双超 13 万亿，与排末的西藏相差 50 余倍。",
                "Horizontal Bar · 水平柱状图", current="ch2")
    slide_chart(prs, "task2_composition.png",
                "T2.2  ·  2024 年三产占比",
                "服务化程度光谱：从新疆 47% 到北京 85%",
                "北京、上海、海南已进入典型服务型经济区间（三产 ≥ 60%）。",
                "100% Stacked Bar · 堆叠条形", current="ch2")
    slide_chart(prs, "task2_dumbbell.png",
                "T2.3  ·  2015 vs 2024 哑铃对比",
                "体量膨胀，位次重排：十年累计增长率对比",
                "西藏 +160% 绝对冠军；头部省份绝对增量远大于尾部。",
                "Dumbbell Plot · 哑铃图", current="ch2")
    slide_chart(prs, "task2_choropleth.png",
                "T2.4  ·  2024 省级 GDP 地图",
                "东南沿海深蓝，西部高原浅淡：2024 省级 GDP 分布",
                "粤苏两省独占深蓝，鲁浙川豫紧随其后。",
                "Choropleth · 分层设色地图", current="ch2")
    slide_chart(prs, "task2_smallmult.png",
                "T2.5  ·  地理网格小多图",
                "地理小多图：31 省十年 GDP 曲线（按方位排列）",
                "橙色 = 头部四强；蓝色 = 其余 27 省；每格右上 = 2024 万亿值。",
                "Small Multiples · 变形地图", current="ch2")
    slide_chart(prs, "task2_h2_clock.png",
                "H2  ·  封面级 HERO 图",
                "结构时钟：几乎所有省份都驶过对角线之上",
                "横纵轴 = 2015/2024 三产占比；偏离对角线 = 服务化幅度。",
                "Clock Chart · 结构时钟", hero=True, current="ch2")

    # ===== Ch.3 相关分析 =====
    slide_chapter_divider(prs, "03", "III",
                          "经济是联动的，\n也是异质的。",
                          "Chapter 03 · 相关分析 · 3 张图",
                          "ch3")
    slide_chart(prs, "task3_parttowhole.png",
                "T3.1  ·  部分-总体 Treemap",
                "全国 GDP 的省际版图：谁贡献了多少？",
                "前 4 强合计 34.9%；前 10 强合计 61.2%。",
                "Treemap · 部分-总体", current="ch3")
    slide_chart(prs, "task3_corr_structure.png",
                "T3.2  ·  结构 vs 增速",
                "结构升级 vs 经济增速：两者有关联吗？",
                "r = −0.21 · 弱负相关。服务化 ≠ 快速增长，追赶靠工业化。",
                "Scatter + Regression · 回归散点", current="ch3")
    slide_chart(prs, "task3_corr_provinces.png",
                "T3.3  ·  省际相关热图",
                "增速的一盘棋：省际年度增速高度同步",
                "31 省两两 Pearson r 均值 +0.79，宏观周期共振显著。",
                "Correlation Heatmap · 相关热图", current="ch3")

    # ===== 结语 =====
    slide_conclusion(prs)

    # ===== 方法论 =====
    slide_colophon(prs)

    out = SLIDES_DIR / "项目展示.pptx"
    prs.save(out)
    print(f"✓ PPT 已生成：{out}  ({len(prs.slides)} 张幻灯片)")


if __name__ == "__main__":
    main()
